import base64
import uuid
from datetime import datetime, timedelta, timezone
from io import BytesIO
from pathlib import Path
from typing import Optional

from fastapi import APIRouter, Depends, HTTPException, BackgroundTasks, status, Form, File, UploadFile
from fastapi.security import OAuth2PasswordBearer
from jose import JWTError, jwt
from pydantic import BaseModel
from sqlalchemy.orm import Session as DBSession

from database import get_db, SessionLocal
from models.student import Student
from models.conversation import Session as ConvSession, Conversation
from services import agent as agent_service
from services import skills as skills_service
from services.memory import summarize_session
from services.usage import check_daily_limit
from services.profile import get_profile_context_for_prompt
from routers.auth import create_access_token, verify_password, hash_password
from config import get_settings

settings = get_settings()
router = APIRouter(prefix="/api/student", tags=["student"])

oauth2_scheme = OAuth2PasswordBearer(tokenUrl="/api/student/login")

# Directory for uploaded session documents
SESSION_DOCS_DIR = Path(__file__).parent.parent / "session_docs"
SESSION_DOCS_DIR.mkdir(exist_ok=True)

ALLOWED_MIME_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}
ALLOWED_EXTENSIONS = {".pdf", ".pptx", ".ppt"}


# --- Request/Response models ---

class LoginRequest(BaseModel):
    name: str
    password: Optional[str] = None


class RegisterRequest(BaseModel):
    name: str
    password: str


class ChangeStudentPasswordRequest(BaseModel):
    current_password: str
    new_password: str


class ChatRequest(BaseModel):
    message: str
    session_id: Optional[int] = None


class ChatResponse(BaseModel):
    reply: str
    session_id: int
    session_mode: Optional[str] = None
    system_prompt: str
    doc_filename: Optional[str] = None
    citations: list[dict] = []


class HistoryMessage(BaseModel):
    id: int
    session_id: int
    role: str
    content: str
    created_at: str
    system_prompt: Optional[str] = None


class SessionOut(BaseModel):
    id: int
    started_at: str
    message_count: int
    last_message: Optional[str] = None
    mode: Optional[str] = None


class ChallengeSessionOut(BaseModel):
    session_id: int
    mode: str
    started_at: str


# --- Auth dependency ---

def get_current_student(token: str = Depends(oauth2_scheme), db: DBSession = Depends(get_db)) -> Student:
    credentials_exception = HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Invalid credentials",
        headers={"WWW-Authenticate": "Bearer"},
    )
    try:
        payload = jwt.decode(token, settings.jwt_secret_key, algorithms=[settings.jwt_algorithm])
        if payload.get("type") != "student":
            raise credentials_exception
        student_id: str = payload.get("sub")
        if student_id is None:
            raise credentials_exception
    except JWTError:
        raise credentials_exception

    student = db.query(Student).filter(Student.id == int(student_id)).first()
    if student is None:
        raise credentials_exception
    return student


# --- Session management ---

def get_or_create_session(db: DBSession, student: Student) -> ConvSession:
    timeout = timedelta(minutes=settings.session_timeout_minutes)
    now = datetime.now(timezone.utc)

    latest = (
        db.query(ConvSession)
        .filter(ConvSession.student_id == student.id, ConvSession.mode == None)  # noqa: E711
        .order_by(ConvSession.started_at.desc())
        .first()
    )

    if latest:
        last_msg = (
            db.query(Conversation)
            .filter(Conversation.session_id == latest.id)
            .order_by(Conversation.created_at.desc())
            .first()
        )
        reference_time = latest.started_at
        if last_msg and last_msg.created_at:
            reference_time = last_msg.created_at

        if reference_time.tzinfo is None:
            reference_time = reference_time.replace(tzinfo=timezone.utc)

        if now - reference_time < timeout and not latest.summarized:
            return latest

    session = ConvSession(student_id=student.id, mode=None)
    db.add(session)
    db.commit()
    db.refresh(session)
    return session


async def trigger_session_summary(student_id: int) -> None:
    db = SessionLocal()
    try:
        timeout = timedelta(minutes=settings.session_timeout_minutes)
        now = datetime.now(timezone.utc)

        sessions = (
            db.query(ConvSession)
            .filter(
                ConvSession.student_id == student_id,
                ConvSession.summarized == False,  # noqa: E712
            )
            .all()
        )

        for session in sessions:
            last_msg = (
                db.query(Conversation)
                .filter(Conversation.session_id == session.id)
                .order_by(Conversation.created_at.desc())
                .first()
            )
            if not last_msg:
                session.summarized = True
                db.commit()
                continue

            ref_time = last_msg.created_at
            if ref_time.tzinfo is None:
                ref_time = ref_time.replace(tzinfo=timezone.utc)

            if now - ref_time > timeout:
                await summarize_session(db, session)
    finally:
        db.close()


# --- Document helpers ---

def _extract_pptx_text(content: bytes) -> str:
    """Extract slide text (titles, body, notes) from a PPTX file."""
    from pptx import Presentation  # type: ignore

    prs = Presentation(BytesIO(content))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[幻灯片 {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
        # Include speaker notes if present
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip() if notes_tf else ""
            if notes_text:
                lines.append(f"  [备注] {notes_text}")
    return "\n".join(lines)


def _save_document(file_bytes: bytes, original_filename: str, content_type: str) -> tuple[Path, str]:
    """
    Save an uploaded file to SESSION_DOCS_DIR.
    Returns (saved_path, media_type).
    PPTX files are converted to plain-text for LLM consumption.
    """
    suffix = Path(original_filename).suffix.lower()

    if suffix not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"不支持的文件类型：{suffix}。请上传 PDF 或 PPTX 文件。")

    unique_stem = uuid.uuid4().hex

    if suffix == ".pdf":
        saved_path = SESSION_DOCS_DIR / f"{unique_stem}.pdf"
        saved_path.write_bytes(file_bytes)
        return saved_path, "application/pdf"
    else:
        # PPTX / PPT — extract text and store as UTF-8 txt
        try:
            text = _extract_pptx_text(file_bytes)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"无法解析 PPTX 文件：{e}")
        saved_path = SESSION_DOCS_DIR / f"{unique_stem}.txt"
        saved_path.write_text(text, encoding="utf-8")
        return saved_path, "text/plain"


# --- Routes ---

@router.get("/sessions", response_model=list[SessionOut])
def list_sessions(
    db: DBSession = Depends(get_db),
    student: Student = Depends(get_current_student),
):
    """Return recent sessions for the current student, latest first."""
    sessions = (
        db.query(ConvSession)
        .filter(ConvSession.student_id == student.id)
        .order_by(ConvSession.started_at.desc())
        .limit(20)
        .all()
    )
    if not sessions:
        return []

    session_ids = [s.id for s in sessions]
    convs = (
        db.query(Conversation)
        .filter(Conversation.session_id.in_(session_ids))
        .order_by(Conversation.created_at.asc())
        .all()
    )

    from collections import defaultdict
    conv_map: dict = defaultdict(list)
    for c in convs:
        conv_map[c.session_id].append(c)

    result = []
    for s in sessions:
        msgs = conv_map[s.id]
        last_user = next((m for m in reversed(msgs) if m.role == "user"), None)
        result.append(SessionOut(
            id=s.id,
            started_at=s.started_at.isoformat(),
            message_count=len(msgs),
            last_message=last_user.content[:50] if last_user else None,
            mode=s.mode,
        ))
    return result


@router.post("/session/new")
def new_session(
    db: DBSession = Depends(get_db),
    student: Student = Depends(get_current_student),
):
    """Explicitly start a fresh conversation session."""
    session = ConvSession(student_id=student.id, mode=None)
    db.add(session)
    db.commit()
    db.refresh(session)
    return {"session_id": session.id}


@router.post("/register")
def student_register(req: RegisterRequest, db: DBSession = Depends(get_db)):
    name = req.name.strip()
    if not name:
        raise HTTPException(status_code=400, detail="姓名不能为空")
    if len(req.password) < 6:
        raise HTTPException(status_code=400, detail="密码至少需要 6 位")

    student = db.query(Student).filter(Student.name == name).first()
    if student:
        if student.hashed_password:
            raise HTTPException(status_code=400, detail="该姓名已注册，请直接登录")
        # Legacy account without password — let them claim it by setting a password
        student.hashed_password = hash_password(req.password)
        db.commit()
    else:
        student = Student(name=name, feishu_user_id=None, hashed_password=hash_password(req.password))
        db.add(student)
        db.commit()
        db.refresh(student)

    token = create_access_token({"sub": str(student.id), "type": "student"})
    return {"access_token": token, "token_type": "bearer", "name": student.name}


@router.post("/login")
def student_login(req: LoginRequest, db: DBSession = Depends(get_db)):
    name = req.name.strip()
    if not name:
        raise HTTPException(status_code=400, detail="姓名不能为空")

    student = db.query(Student).filter(Student.name == name).first()
    if not student:
        raise HTTPException(status_code=400, detail="该姓名未注册，请先注册")

    if student.hashed_password:
        if not req.password or not verify_password(req.password, student.hashed_password):
            raise HTTPException(status_code=400, detail="密码错误")
    # Legacy accounts without a password: allow login without password

    token = create_access_token({"sub": str(student.id), "type": "student"})
    return {"access_token": token, "token_type": "bearer", "name": student.name}


@router.post("/change-password")
def student_change_password(
    req: ChangeStudentPasswordRequest,
    student: Student = Depends(get_current_student),
    db: DBSession = Depends(get_db),
):
    if student.hashed_password and not verify_password(req.current_password, student.hashed_password):
        raise HTTPException(status_code=400, detail="当前密码不正确")
    if len(req.new_password) < 6:
        raise HTTPException(status_code=400, detail="新密码至少需要 6 位")
    student.hashed_password = hash_password(req.new_password)
    db.commit()
    return {"ok": True}


@router.post("/chat", response_model=ChatResponse)
async def student_chat(
    req: ChatRequest,
    background_tasks: BackgroundTasks,
    db: DBSession = Depends(get_db),
    student: Student = Depends(get_current_student),
):
    text = req.message.strip()
    if not text:
        raise HTTPException(status_code=400, detail="Message cannot be empty")

    if req.session_id is not None:
        session = db.query(ConvSession).filter(
            ConvSession.id == req.session_id,
            ConvSession.student_id == student.id,
        ).first()
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
    else:
        session = get_or_create_session(db, student)

    check_daily_limit(db, student)
    reply, input_tokens, output_tokens, system_prompt, citations = agent_service.chat(db, student, session, text)

    db.add(Conversation(student_id=student.id, session_id=session.id, role="user", content=text))
    db.add(Conversation(
        student_id=student.id, session_id=session.id,
        role="assistant", content=reply,
        input_tokens=input_tokens, output_tokens=output_tokens,
        system_prompt=system_prompt,
    ))
    db.commit()

    background_tasks.add_task(trigger_session_summary, student.id)

    doc_filename = None
    if session.doc_path:
        # Return the original media type hint as a simple label
        doc_filename = Path(session.doc_path).name

    return {"reply": reply, "session_id": session.id, "session_mode": session.mode, "system_prompt": system_prompt, "doc_filename": doc_filename, "citations": citations}


@router.post("/chat/stream")
async def student_chat_stream(
    req: ChatRequest,
    background_tasks: BackgroundTasks,
    db: DBSession = Depends(get_db),
    student: Student = Depends(get_current_student),
):
    """SSE streaming variant of /chat."""
    import json as _json
    from starlette.responses import StreamingResponse

    text = req.message.strip()
    if not text:
        raise HTTPException(status_code=400, detail="Message cannot be empty")

    if req.session_id is not None:
        session = db.query(ConvSession).filter(
            ConvSession.id == req.session_id,
            ConvSession.student_id == student.id,
        ).first()
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
    else:
        session = get_or_create_session(db, student)

    check_daily_limit(db, student)

    session_id = session.id
    student_id = student.id

    def generate():
        full_reply = []
        system_prompt = ""
        citations = []
        input_tokens = 0
        output_tokens = 0

        for kind, data in agent_service.chat_stream(db, student, session, text):
            if kind == "setup":
                system_prompt = data["system_prompt"]
                citations = data["citations"]
            elif kind == "delta":
                full_reply.append(data)
                yield f"data: {_json.dumps({'t': data}, ensure_ascii=False)}\n\n"
            elif kind == "done":
                input_tokens = data.get("input_tokens", 0)
                output_tokens = data.get("output_tokens", 0)

        # Send final event with metadata
        yield f"data: {_json.dumps({'done': True, 'session_id': session_id, 'system_prompt': system_prompt, 'citations': citations}, ensure_ascii=False)}\n\n"

        # Record to DB
        reply = "".join(full_reply)
        db.add(Conversation(student_id=student_id, session_id=session_id, role="user", content=text))
        db.add(Conversation(
            student_id=student_id, session_id=session_id,
            role="assistant", content=reply,
            input_tokens=input_tokens, output_tokens=output_tokens,
            system_prompt=system_prompt,
        ))
        db.commit()

        background_tasks.add_task(trigger_session_summary, student_id)

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.post("/chat/upload", response_model=ChatResponse)
async def student_chat_upload(
    background_tasks: BackgroundTasks,
    message: str = Form(...),
    file: UploadFile = File(...),
    session_id: Optional[int] = Form(None),
    db: DBSession = Depends(get_db),
    student: Student = Depends(get_current_student),
):
    """Send a message with a document (PDF or PPTX). The document is attached to the session
    and re-sent to the LLM on every subsequent turn in the same session."""
    text = message.strip()
    if not text:
        raise HTTPException(status_code=400, detail="Message cannot be empty")

    file_bytes = await file.read()
    if len(file_bytes) > 20 * 1024 * 1024:  # 20 MB limit
        raise HTTPException(status_code=413, detail="文件过大，最大支持 20 MB。")

    saved_path, media_type = _save_document(file_bytes, file.filename or "upload", file.content_type or "")

    # Resolve / create session
    if session_id is not None:
        session = db.query(ConvSession).filter(
            ConvSession.id == session_id,
            ConvSession.student_id == student.id,
        ).first()
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
    else:
        session = get_or_create_session(db, student)

    # If this session already had a different document, delete the old file
    if session.doc_path and session.doc_path != str(saved_path):
        old = Path(session.doc_path)
        if old.exists():
            old.unlink(missing_ok=True)

    # Attach document to session
    session.doc_path = str(saved_path)
    session.doc_media_type = media_type
    db.commit()

    check_daily_limit(db, student)

    # Build the document dict to pass for this first call
    document = {
        "media_type": media_type,
        "data": base64.b64encode(saved_path.read_bytes()).decode(),
    }
    reply, input_tokens, output_tokens, system_prompt, citations = agent_service.chat(
        db, student, session, text, document=document
    )

    db.add(Conversation(student_id=student.id, session_id=session.id, role="user", content=text))
    db.add(Conversation(
        student_id=student.id, session_id=session.id,
        role="assistant", content=reply,
        input_tokens=input_tokens, output_tokens=output_tokens,
        system_prompt=system_prompt,
    ))
    db.commit()

    background_tasks.add_task(trigger_session_summary, student.id)

    return {
        "reply": reply,
        "session_id": session.id,
        "session_mode": session.mode,
        "system_prompt": system_prompt,
        "doc_filename": file.filename,
        "citations": citations,
    }


@router.get("/history", response_model=list[HistoryMessage])
def student_history(
    session_id: Optional[int] = None,
    db: DBSession = Depends(get_db),
    student: Student = Depends(get_current_student),
):
    q = db.query(Conversation).filter(Conversation.student_id == student.id)
    if session_id is not None:
        q = q.filter(Conversation.session_id == session_id)
    msgs = (
        q.order_by(Conversation.created_at.desc(), Conversation.id.desc())
        .limit(60)
        .all()
    )
    msgs.reverse()
    return [
        {
            "id": m.id,
            "session_id": m.session_id,
            "role": m.role,
            "content": m.content,
            "created_at": m.created_at.isoformat() if m.created_at else "",
            "system_prompt": m.system_prompt,
        }
        for m in msgs
    ]


@router.post("/challenge/start", response_model=ChallengeSessionOut)
def challenge_start(
    db: DBSession = Depends(get_db),
    student: Student = Depends(get_current_student),
):
    """Start or resume a challenge session. Returns existing active challenge if one exists."""
    existing = (
        db.query(ConvSession)
        .filter(
            ConvSession.student_id == student.id,
            ConvSession.mode == "challenge",
            ConvSession.summarized == False,  # noqa: E712
        )
        .order_by(ConvSession.started_at.desc())
        .first()
    )
    if existing:
        return {
            "session_id": existing.id,
            "mode": existing.mode,
            "started_at": existing.started_at.isoformat(),
        }

    session = ConvSession(student_id=student.id, mode="challenge")
    db.add(session)
    db.commit()
    db.refresh(session)

    # Generate an opening greeting from the AI for new sessions
    # Build KP skill info (name + description)
    kp_skills = [s for s in skills_service.list_skills() if s["enabled"] and s["type"] == "knowledge_point"]
    if kp_skills:
        topics_info = "\n".join(
            f"- {s['name']}" + (f"：{s['description']}" if s.get("description") else "")
            for s in kp_skills
        )
    else:
        topics_info = "（暂无预设知识点）"

    # Get student profile context
    profile_context = get_profile_context_for_prompt(student.id)

    # Get recent normal session conversations (last 10 messages)
    recent_convs = (
        db.query(Conversation)
        .join(ConvSession, Conversation.session_id == ConvSession.id)
        .filter(
            ConvSession.student_id == student.id,
            ConvSession.mode.is_(None),
        )
        .order_by(Conversation.created_at.desc())
        .limit(10)
        .all()
    )
    if recent_convs:
        recent_msgs = "\n".join(
            f"[{c.role}]: {c.content[:200]}" for c in reversed(recent_convs)
        )
    else:
        recent_msgs = "（无近期对话记录）"

    kickoff = f"""（系统提示：挑战模式新会话已启动。

**学生画像**：
{profile_context if profile_context else '暂无画像数据'}

**近期对话记录（最近10条）**：
{recent_msgs}

**可供出题的预设知识领域**：
{topics_info}

请根据以上信息，用中文向学生打招呼，简要说明挑战模式的玩法，然后**主动推荐 1-2 个基于学生薄弱环节或近期话题的挑战方向**，并说明理由。同时告知学生可以接受推荐、选择其他知识领域，或自由输入任何想挑战的主题——即使没有预设知识领域也可以继续挑战。）"""
    reply, input_tokens, output_tokens, system_prompt, _citations = agent_service.chat(db, student, session, kickoff)
    db.add(Conversation(
        student_id=student.id, session_id=session.id,
        role="assistant", content=reply,
        input_tokens=input_tokens, output_tokens=output_tokens,
        system_prompt=system_prompt,
    ))
    db.commit()

    return {
        "session_id": session.id,
        "mode": session.mode,
        "started_at": session.started_at.isoformat(),
    }


@router.get("/challenge/active", response_model=Optional[ChallengeSessionOut])
def challenge_active(
    db: DBSession = Depends(get_db),
    student: Student = Depends(get_current_student),
):
    """Check if there is an ongoing (not summarized) challenge session."""
    existing = (
        db.query(ConvSession)
        .filter(
            ConvSession.student_id == student.id,
            ConvSession.mode == "challenge",
            ConvSession.summarized == False,  # noqa: E712
        )
        .order_by(ConvSession.started_at.desc())
        .first()
    )
    if not existing:
        return None
    return {
        "session_id": existing.id,
        "mode": existing.mode,
        "started_at": existing.started_at.isoformat(),
    }


# --- Onboarding ("学习初心") mode ---

@router.post("/onboarding/start", response_model=ChallengeSessionOut)
def onboarding_start(
    db: DBSession = Depends(get_db),
    student: Student = Depends(get_current_student),
):
    """Start or resume an onboarding session to discover learning motivation."""
    existing = (
        db.query(ConvSession)
        .filter(
            ConvSession.student_id == student.id,
            ConvSession.mode == "onboarding",
            ConvSession.summarized == False,  # noqa: E712
        )
        .order_by(ConvSession.started_at.desc())
        .first()
    )
    if existing:
        return {
            "session_id": existing.id,
            "mode": existing.mode,
            "started_at": existing.started_at.isoformat(),
        }

    session = ConvSession(student_id=student.id, mode="onboarding")
    db.add(session)
    db.commit()
    db.refresh(session)

    # Get student profile context (may include previous learning-motivation aspect)
    profile_context = get_profile_context_for_prompt(student.id)

    # Gather previous onboarding conversation history
    prev_onboarding_convs = (
        db.query(Conversation)
        .join(ConvSession, Conversation.session_id == ConvSession.id)
        .filter(
            ConvSession.student_id == student.id,
            ConvSession.mode == "onboarding",
            ConvSession.id != session.id,  # exclude current new session
        )
        .order_by(Conversation.created_at.desc())
        .limit(20)
        .all()
    )
    if prev_onboarding_convs:
        prev_msgs = "\n".join(
            f"[{c.role}]: {c.content[:300]}" for c in reversed(prev_onboarding_convs)
        )
    else:
        prev_msgs = ""

    profile_line = ('- 已有画像：\n' + profile_context) if profile_context else '- 暂无画像数据（新学生）'

    if prev_msgs:
        kickoff = f"""（系统提示：学习初心对话已启动。

**学生信息**：
- 姓名：{student.name}
{profile_line}

**之前的初心对话记录**：
{prev_msgs}

请用中文向学生打招呼。这不是第一次初心对话了，请结合上面之前的初心聊天记录，继续深入了解学生的**学习动机**和对操作系统课程的**期望**。你可以从之前的回答中找到值得深入的点追问下去，了解他/她的期望或目标是否有变化，或从新的角度切入挖掘更深层的想法。不要重复之前已经聊过的内容。

**注意**：请保持轻松友好，不要一上来就列出很多问题。提出一个有针对性的问题即可。）"""
    else:
        kickoff = f"""（系统提示：学习初心对话已启动。

**学生信息**：
- 姓名：{student.name}
{profile_line}

这是学生第一次初心对话。请用中文向学生打招呼，自我介绍你是课程AI助教，简要说明你想花几分钟了解一下他/她学习这门课的想法。然后从一个轻松的开放性问题开始，比如了解学生的背景或为什么选了这门课。最终目标是了解学生**期望从操作系统这门课中获得什么**。

**注意**：请保持轻松友好，不要一上来就列出很多问题，每次只问一个。）"""

    reply, input_tokens, output_tokens, system_prompt, _citations = agent_service.chat(db, student, session, kickoff)
    db.add(Conversation(
        student_id=student.id, session_id=session.id,
        role="assistant", content=reply,
        input_tokens=input_tokens, output_tokens=output_tokens,
        system_prompt=system_prompt,
    ))
    db.commit()

    return {
        "session_id": session.id,
        "mode": session.mode,
        "started_at": session.started_at.isoformat(),
    }



@router.get("/onboarding/active", response_model=Optional[ChallengeSessionOut])
def onboarding_active(
    db: DBSession = Depends(get_db),
    student: Student = Depends(get_current_student),
):
    """Check if there is an ongoing (not summarized) onboarding session."""
    existing = (
        db.query(ConvSession)
        .filter(
            ConvSession.student_id == student.id,
            ConvSession.mode == "onboarding",
            ConvSession.summarized == False,  # noqa: E712
        )
        .order_by(ConvSession.started_at.desc())
        .first()
    )
    if not existing:
        return None
    return {
        "session_id": existing.id,
        "mode": existing.mode,
        "started_at": existing.started_at.isoformat(),
    }

